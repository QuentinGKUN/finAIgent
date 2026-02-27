from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
IMG_DIR = ROOT / "图片"
OUT_FILE = ROOT / f"财务分析智能助手_系统总结_{datetime.now().strftime('%Y%m%d')}.pptx"


def set_run_style(run, size=Pt(20), bold=False, color=RGBColor(15, 23, 42), font="Microsoft YaHei"):
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_style(r, size=Pt(34), bold=True, color=RGBColor(2, 6, 23))
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(11.5), Inches(0.6))
        sp = sub_box.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        set_run_style(sr, size=Pt(16), color=RGBColor(71, 85, 105))


def add_bullets(slide, items: Iterable[str], left=0.9, top=1.9, width=11.9, height=4.9, level=0):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.text = item
        p.space_after = Pt(10)
        p.space_before = Pt(2)
        for run in p.runs:
            set_run_style(run, size=Pt(20), color=RGBColor(30, 41, 59))


def add_box(slide, x, y, w, h, text, fill=RGBColor(240, 249, 255), line=RGBColor(14, 116, 144)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run_style(r, size=Pt(16), bold=True, color=RGBColor(8, 47, 73))
    return shape


def connect(slide, a, b, color=RGBColor(51, 65, 85)):
    ax = a.left + a.width
    ay = a.top + a.height // 2
    bx = b.left
    by = b.top + b.height // 2
    line = slide.shapes.add_connector(1, ax, ay, bx, by)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)


def fit_picture(slide, img_path: Path, left, top, width, height):
    from PIL import Image

    img = Image.open(img_path)
    iw, ih = img.size
    img_ratio = iw / ih
    box_ratio = width / height
    if img_ratio > box_ratio:
        target_h = width / img_ratio
        y = top + (height - target_h) / 2
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(y), width=Inches(width))
    else:
        target_w = height * img_ratio
        x = left + (width - target_w) / 2
        slide.shapes.add_picture(str(img_path), Inches(x), Inches(top), height=Inches(height))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    s = prs.slides.add_slide(blank)
    add_title(s, "财务分析智能助手（Champion）系统总结", "系统实现 · Tools 体系 · AI 架构 · 运行完整性")
    date_box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.0), Inches(0.7))
    p = date_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"汇报日期：{datetime.now().strftime('%Y-%m-%d')}"
    set_run_style(r, size=Pt(18), color=RGBColor(51, 65, 85))
    highlight = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.9), Inches(11.8), Inches(3.8))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(248, 250, 252)
    highlight.line.color.rgb = RGBColor(203, 213, 225)
    tf = highlight.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "目标：展示该项目从前端交互到后端智能路由、数据工具调用与报告导出的完整工程闭环。"
    for run in p.runs:
        set_run_style(run, size=Pt(20), bold=True, color=RGBColor(15, 23, 42))
    p2 = tf.add_paragraph()
    p2.text = "关键亮点：多数据源融合（AkShare/EODHD/ValueCell）+ 模式化 AI 路由 + 可视化图表与可编辑报告。"
    for run in p2.runs:
        set_run_style(run, size=Pt(18), color=RGBColor(30, 41, 59))

    # 2. Project scope
    s = prs.slides.add_slide(blank)
    add_title(s, "1. 项目实现范围", "端到端财务研究工作台")
    add_bullets(
        s,
        [
            "统一对话入口：支持普通问答、数据对比、深度分析三种交互路径。",
            "数据能力：A 股（AkShare）、美股/港股（EODHD）、深度分析（ValueCell）。",
            "AI 能力：意图判断、公司识别、多轮记忆、证据融合、结论生成。",
            "可视化能力：自动图表（Chart.js）+ 数据表格抽取 + 引用来源追溯。",
            "交付能力：会话历史管理 + 导出研究报告（可勾选/编辑/打印）。",
        ],
    )

    # 3. System architecture
    s = prs.slides.add_slide(blank)
    add_title(s, "2. 系统总体架构", "前端 / 后端 / 数据与工具 / AI 服务分层")
    fe = add_box(s, 0.8, 1.8, 2.6, 1.0, "React 前端\nweb/src/App.jsx", RGBColor(224, 242, 254))
    api = add_box(s, 3.8, 1.8, 2.8, 1.0, "Go API 网关\n/api/chat", RGBColor(219, 234, 254))
    mem = add_box(s, 7.0, 1.8, 2.5, 1.0, "会话与记忆\nSQLite / in-memory", RGBColor(226, 232, 240))
    llm = add_box(s, 9.9, 1.8, 2.5, 1.0, "GLM LLM\n总结/规划", RGBColor(255, 237, 213), RGBColor(194, 65, 12))

    tool1 = add_box(s, 3.8, 3.3, 2.8, 1.0, "AkShare Bridge\nA股财务/榜单", RGBColor(209, 250, 229), RGBColor(5, 150, 105))
    tool2 = add_box(s, 7.0, 3.3, 2.5, 1.0, "EODHD\n美股/港股财务", RGBColor(209, 250, 229), RGBColor(5, 150, 105))
    tool3 = add_box(s, 9.9, 3.3, 2.5, 1.0, "ValueCell\n深度研究", RGBColor(209, 250, 229), RGBColor(5, 150, 105))

    report = add_box(s, 3.8, 4.9, 8.6, 1.1, "报告引擎 /api/report ：结构化问答 + 图表 + 引用 + 打印输出", RGBColor(243, 232, 255), RGBColor(109, 40, 217))

    connect(s, fe, api)
    connect(s, api, mem)
    connect(s, mem, llm)
    connect(s, api, tool1)
    connect(s, api, tool2)
    connect(s, api, tool3)
    connect(s, api, report)

    # 4. AI architecture
    s = prs.slides.add_slide(blank)
    add_title(s, "3. AI 架构设计", "模式控制 + 数据路由 + 证据汇聚")
    add_bullets(
        s,
        [
            "模式层：`normal`（模型问答）/ `#mode:pro`（数据对比）/ `#mode:deep`（深度分析）。",
            "路由层：先判断问题是否需要数据，再决定走 Compare / Finance / Rank / Board / Deep。",
            "记忆层：会话摘要 + 历史证据合并，支持追问场景上下文连续性。",
            "执行层：优先工具取数（AkShare/EODHD/ValueCell），失败时自动降级回模型。",
            "生成层：LLM 基于“当前证据 + 历史证据”生成结论，并限制无证据扩写。",
            "可观测层：进度状态、数据通道标识、引用链路、查询记录可回放。",
        ],
        top=1.9,
        height=5.1,
    )

    # 5. Tools landscape
    s = prs.slides.add_slide(blank)
    add_title(s, "4. Tools 体系与职责", "工具协作驱动可核验分析")
    rows, cols = 7, 4
    table = s.shapes.add_table(rows, cols, Inches(0.7), Inches(1.6), Inches(12.0), Inches(4.9)).table
    headers = ["Tool/组件", "职责", "触发场景", "输出"]
    data = [
        ["AkShare", "A股财务、板块、个股排行", "数据对比/榜单查询", "结构化时序数据 + 来源链接"],
        ["EODHD", "美股/港股 fundamentals", "海外公司财务问题", "财务指标序列 + 官方文档引用"],
        ["ValueCell", "深度分析、代理执行", "深度分析模式", "研究结论 + 工具调用痕迹"],
        ["GLM", "意图判断与结论生成", "所有模式总结阶段", "自然语言分析 + 风险提示"],
        ["Chart.js", "图表渲染", "检测到趋势/时间序列", "折线/柱状图可视化"],
        ["Report Engine", "会话转投研报告", "导出报告", "可编辑条目 + 打印/PDF"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
    for r_idx, row in enumerate(data, start=1):
        for c_idx, val in enumerate(row):
            table.cell(r_idx, c_idx).text = val
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            tf = cell.text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    set_run_style(run, size=Pt(13), bold=(r == 0), color=RGBColor(15, 23, 42))
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(226, 232, 240)

    # 6. Sequence
    s = prs.slides.add_slide(blank)
    add_title(s, "5. 核心调用链路（深度分析示例）", "体现系统架构完整性")
    add_bullets(
        s,
        [
            "① 用户在前端点击“深度分析”并提交问题。",
            "② Go 后端接收请求，进入 `#mode:deep` 分支，提取目标公司与上下文记忆。",
            "③ 调用 ValueCell bridge 执行研究流程（并自动处理 Provider 配置与重试）。",
            "④ 汇总 ValueCell 返回证据，追加引用（查询记录、官方仓库）并落库。",
            "⑤ GLM 对证据进行结构化总结，生成可读结论。",
            "⑥ 前端渲染数据通道徽标、引用区、图表区，保证可核验与可追溯。",
            "⑦ 导出报告页可二次编辑、勾选删除、打印输出，形成投研初稿。",
        ],
        top=1.9,
        height=5.1,
    )

    # 7. Engineering quality
    s = prs.slides.add_slide(blank)
    add_title(s, "6. 工程完整性与高级能力", "稳定性、可观测性、可扩展性")
    add_bullets(
        s,
        [
            "稳定性：多工具失败回退（数据查询失败 -> 模型回退），避免中断对话。",
            "可观测性：聊天进度阶段（memory/route/query/summary）+ 耗时展示。",
            "可追溯性：每次工具调用均可打开查询记录，报告中保留引用来源。",
            "可扩展性：模式路由和桥接脚本解耦，后续可接入更多金融数据服务。",
            "可运营性：Windows 一键启动脚本、端口分层（5173/3000/8010）部署清晰。",
        ],
    )

    # 8+. Screenshots from image directory
    images = sorted([p for p in IMG_DIR.glob("*") if p.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}])
    if images:
        per_slide = 2
        total = (len(images) + per_slide - 1) // per_slide
        for idx in range(total):
            s = prs.slides.add_slide(blank)
            add_title(s, f"系统截图集（{idx+1}/{total}）", "素材来源：项目目录 /图片")
            batch = images[idx * per_slide : (idx + 1) * per_slide]
            if len(batch) == 1:
                fit_picture(s, batch[0], 0.9, 1.6, 11.7, 5.4)
                cap = s.shapes.add_textbox(Inches(0.9), Inches(7.0 - 0.35), Inches(11.7), Inches(0.3))
                cp = cap.text_frame.paragraphs[0]
                cr = cp.add_run()
                cr.text = batch[0].name
                set_run_style(cr, size=Pt(11), color=RGBColor(71, 85, 105))
            else:
                fit_picture(s, batch[0], 0.8, 1.6, 5.9, 4.8)
                fit_picture(s, batch[1], 6.6, 1.6, 5.9, 4.8)
                cap1 = s.shapes.add_textbox(Inches(0.8), Inches(6.45), Inches(5.9), Inches(0.3))
                cap2 = s.shapes.add_textbox(Inches(6.6), Inches(6.45), Inches(5.9), Inches(0.3))
                p1 = cap1.text_frame.paragraphs[0]
                r1 = p1.add_run()
                r1.text = batch[0].name
                set_run_style(r1, size=Pt(11), color=RGBColor(71, 85, 105))
                p2 = cap2.text_frame.paragraphs[0]
                r2 = p2.add_run()
                r2.text = batch[1].name
                set_run_style(r2, size=Pt(11), color=RGBColor(71, 85, 105))

    # Final slide
    s = prs.slides.add_slide(blank)
    add_title(s, "总结", "从可用原型升级到专业投研工作台")
    add_bullets(
        s,
        [
            "系统已形成“对话 -> 数据工具 -> AI 总结 -> 图表/报告”闭环架构。",
            "通过模式化路由将普通问答、数据对比、深度分析统一到一个交互入口。",
            "通过引用与查询记录保证分析结果可核验，满足专业场景信任要求。",
            "后续可增强：ResearchAgent 直连稳定性、策略因子模块、报告模板化输出。",
        ],
        top=2.0,
    )

    prs.save(str(OUT_FILE))
    print(str(OUT_FILE))


if __name__ == "__main__":
    main()
